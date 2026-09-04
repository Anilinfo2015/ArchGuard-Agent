"""Generate the internal pitch deck as a real .pptx with native diagrams.

Draws the flow, sequence, and role diagrams as PowerPoint shapes (no external image
renderer needed) so the deck is fully editable. Run:

    python pitch/build_deck.py

Produces pitch/ArchGuard-Pitch.pptx.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

BLUE = RGBColor(0x1F, 0x4E, 0x79)
LIGHT = RGBColor(0xDA, 0xE3, 0xF3)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
GREY = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, left, top, width, height, text, size=18, bold=False,
            color=GREY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def bullets(s, left, top, width, height, items, size=18):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = GREY
    return box


def box(s, left, top, width, height, text, fill=LIGHT, line=BLUE,
        font=BLUE, size=13, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(1.5)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = font
    return sp


def connect(s, a, b, color=BLUE):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1])
    c.line.color.rgb = color
    c.line.width = Pt(1.75)
    line = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = line.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    line.append(tail)
    return c


def header(s, title, subtitle=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "   " + title
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        textbox(s, 0.3, 1.15, 12.7, 0.5, subtitle, size=16, color=GREY, bold=True)


# 1. Title
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
textbox(s, 1, 2.4, 11.3, 1.2, "ArchGuard-AI", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, 1, 3.7, 11.3, 0.9, "Continuous architecture review at every step — not once",
        size=26, color=LIGHT, align=PP_ALIGN.CENTER)
textbox(s, 1, 5.4, 11.3, 0.6, "Internal hackathon pitch", size=16, color=LIGHT, align=PP_ALIGN.CENTER)

# 2. Problem
s = slide()
header(s, "The problem")
bullets(s, 0.8, 1.6, 11.7, 5, [
    "Architecture is reviewed once at an ARB, then drifts silently.",
    "Rules live in Confluence and people's heads — not enforced in code.",
    "By the time drift is found, it is expensive to unwind.",
    "Generic AI PR reviewers are non-deterministic and cannot gate a merge.",
], size=22)
textbox(s, 0.8, 6.0, 11.7, 0.8,
        "\u201cWe review architecture at the start, then hope. Drift is discovered late.\u201d",
        size=20, bold=True, color=BLUE)

# 3. The idea
s = slide()
header(s, "The idea")
bullets(s, 0.8, 1.6, 11.7, 4.5, [
    "Architects write rules in plain English (org / domain / team / service, HLD + LLD).",
    "An LLM compiles each rule into architecture-as-code — once, at authoring time.",
    "The pipeline runs the compiled rules deterministically on every PR.",
    "Findings carry a policy ID + file:line. High-confidence drift blocks the merge.",
], size=22)
textbox(s, 0.8, 6.1, 11.7, 0.8,
        "The model proposes at authoring time. The compiled rule decides at runtime.",
        size=20, bold=True, color=GREEN)

# 4. Flow diagram
s = slide()
header(s, "Flow diagram")
textbox(s, 0.5, 1.3, 6, 0.4, "Author once (LLM)", size=15, bold=True, color=BLUE)
b1 = box(s, 0.5, 1.8, 2.6, 1.0, "Rules in English\n(MD via PR / UI)")
b2 = box(s, 3.5, 1.8, 2.4, 1.0, "LLM compiler\nskill/")
b3 = box(s, 6.3, 1.8, 2.6, 1.0, "Architecture-as-code\npolicy pack", fill=RGBColor(0xD5, 0xE8, 0xD4), line=GREEN, font=GREEN)
connect(s, (Inches(3.1), Inches(2.3)), (Inches(3.5), Inches(2.3)))
connect(s, (Inches(5.9), Inches(2.3)), (Inches(6.3), Inches(2.3)))
textbox(s, 0.5, 3.3, 8, 0.4, "Enforce every PR (deterministic)", size=15, bold=True, color=BLUE)
p1 = box(s, 0.5, 3.8, 2.2, 1.0, "Pull request")
p2 = box(s, 3.2, 3.8, 3.0, 1.0, "Native parsers\nStructurizr / dep-cruiser / TF")
p3 = box(s, 6.7, 3.8, 2.4, 1.0, "Engine evaluates\nengine/")
d1 = box(s, 9.6, 3.4, 3.0, 0.9, "Blocking? -> PR fails\npolicyId + file:line", fill=RGBColor(0xF8, 0xCE, 0xCE), line=RGBColor(0xC0, 0x39, 0x2B), font=RGBColor(0xC0, 0x39, 0x2B))
d2 = box(s, 9.6, 4.6, 3.0, 0.9, "Otherwise -> advisory", fill=RGBColor(0xFF, 0xF2, 0xCC), line=AMBER, font=AMBER)
connect(s, (Inches(2.7), Inches(4.3)), (Inches(3.2), Inches(4.3)))
connect(s, (Inches(6.2), Inches(4.3)), (Inches(6.7), Inches(4.3)))
connect(s, (Inches(9.1), Inches(4.2)), (Inches(9.6), Inches(3.85)), color=RGBColor(0xC0, 0x39, 0x2B))
connect(s, (Inches(9.1), Inches(4.4)), (Inches(9.6), Inches(5.05)), color=AMBER)
connect(s, (Inches(7.6), Inches(1.8)), (Inches(7.9), Inches(3.8)), color=GREEN)

# 5. Sequence diagram
s = slide()
header(s, "Sequence diagram")
actors = ["Architect", "UI / PR", "Skill (LLM)", "Repo (main)", "Developer", "Pipeline", "Engine"]
n = len(actors)
left0, span, top = 0.5, 12.3, 1.5
xs = [left0 + span * (i + 0.5) / n for i in range(n)]
for i, name in enumerate(actors):
    box(s, xs[i] - 0.75, top, 1.5, 0.6, name, size=11)
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xs[i]), Inches(top + 0.6), Inches(xs[i]), Inches(6.9))
    ln.line.color.rgb = GREY
    ln.line.width = Pt(1)


def msg(a, b, y, label):
    connect(s, (Inches(xs[a]), Inches(y)), (Inches(xs[b]), Inches(y)), color=BLUE)
    lo = min(xs[a], xs[b])
    textbox(s, lo, y - 0.32, abs(xs[b] - xs[a]) + 0.2, 0.3, label, size=10, color=GREY,
            align=PP_ALIGN.CENTER if b > a else PP_ALIGN.CENTER)


steps = [
    (0, 1, "write rule"), (1, 2, "submit"), (2, 2, "compile to code"),
    (2, 3, "PR policy pack"), (0, 3, "review & merge"), (4, 5, "open PR"),
    (5, 5, "run parsers"), (5, 6, "evaluate"), (6, 5, "findings"), (5, 4, "block / advise"),
]
y = 2.6
for a, b, label in steps:
    if a == b:
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xs[a]), Inches(y - 0.1), Inches(0.9), Inches(0.28))
        sp.fill.solid(); sp.fill.fore_color.rgb = LIGHT; sp.line.color.rgb = BLUE
        textbox(s, xs[a] + 0.05, y - 0.12, 1.6, 0.3, label, size=9, color=BLUE)
    else:
        msg(a, b, y, label)
    y += 0.42

# 6. Role diagram
s = slide()
header(s, "Role diagram")
box(s, 0.7, 1.7, 3.0, 1.6, "Architect\n\nAuthors org/team rules\nReviews compiled code", size=13)
box(s, 5.1, 1.7, 3.0, 1.6, "Skill / LLM\n(authoring time)\n\nCompiles English -> predicates", size=13, fill=RGBColor(0xD5, 0xE8, 0xD4), line=GREEN, font=GREEN)
box(s, 9.5, 1.7, 3.0, 1.6, "Developer\n\nOpens PR\nFixes flagged drift", size=13)
box(s, 5.1, 4.6, 3.0, 1.4, "Pipeline\n(runtime)\n\nParsers + engine on every PR", size=13, fill=RGBColor(0xFF, 0xF2, 0xCC), line=AMBER, font=AMBER)
connect(s, (Inches(3.7), Inches(2.5)), (Inches(5.1), Inches(2.5)))
connect(s, (Inches(8.1), Inches(2.5)), (Inches(9.5), Inches(2.5)))
connect(s, (Inches(6.6), Inches(3.3)), (Inches(6.6), Inches(4.6)), color=GREEN)
connect(s, (Inches(11.0), Inches(3.3)), (Inches(8.1), Inches(4.9)))

# 7. Architecture-as-code example
s = slide()
header(s, "What architecture-as-code looks like")
textbox(s, 0.8, 1.4, 11.7, 0.8, "English rule:", size=18, bold=True, color=BLUE)
textbox(s, 0.8, 1.9, 11.7, 0.8,
        "\u201cModules in payments must not call legacy-billing synchronously.\u201d",
        size=18, color=GREY)
code = ("{\n"
        '  "policyId": "ARC-COM-003",\n'
        '  "scope": "domain:payments",\n'
        '  "tier": "lld",\n'
        '  "severity": "block",\n'
        '  "predicate": "require_async_cross_context(target=\'legacy-billing\')",\n'
        '  "source": "rules/payments.md#rule-1"\n'
        "}")
cb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.9), Inches(11.7), Inches(3.4))
cb.fill.solid(); cb.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E); cb.line.fill.background()
tf = cb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = code
r.font.size = Pt(16); r.font.name = "Consolas"; r.font.color.rgb = RGBColor(0x9C, 0xDC, 0xFE)

# 8. Parallel work
s = slide()
header(s, "Built for parallel work")
rows = [
    ("Skill / Authoring (LLM)", "skill/"),
    ("Engine (deterministic)", "engine/"),
    ("Pipeline (CI trigger)", "pipeline/"),
    ("UI (author + review)", "ui/"),
    ("Pitch", "pitch/"),
]
y = 1.8
for area, folder in rows:
    box(s, 0.8, y, 6.5, 0.7, area, size=15, fill=LIGHT)
    box(s, 7.5, y, 3.0, 0.7, folder, size=15, fill=WHITE, font=GREY, line=GREY)
    box(s, 10.7, y, 1.8, 0.7, "owner", size=13, fill=WHITE, font=GREY, line=GREY)
    y += 0.85
textbox(s, 0.8, 6.3, 11.7, 0.6, "Four independent workstreams, one JSON contract between them.",
        size=18, bold=True, color=BLUE)

# 9. Why it wins + ask
s = slide()
header(s, "Why it wins & the ask")
bullets(s, 0.8, 1.5, 11.7, 3.2, [
    "Deterministic gate — reproducible and auditable, unlike generic AI reviewers.",
    "Evidence-bound — every finding has a policy ID and file:line.",
    "Author once, enforce always — LLM effort spent at authoring, not per PR.",
    "Continuous — architecture review at every step, seamless for developers.",
], size=20)
textbox(s, 0.8, 4.9, 11.7, 1.6,
        "The ask:\n"
        "• Pick your area and grab it in CODEOWNERS.\n"
        "• Demo: English rule -> compiled policy -> blocked PR, end to end.\n"
        "• Stretch: live UI authoring + advisory comments on a real PR.",
        size=18, bold=False, color=GREEN)

out = Path(__file__).parent / "ArchGuard-Pitch.pptx"
prs.save(str(out))
print(f"wrote {out}")
